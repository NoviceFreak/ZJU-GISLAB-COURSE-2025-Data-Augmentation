from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os
import sys


def pdf_to_pptx_enhanced(pdf_path, output_path, dpi=200):
    """
    Enhanced PDF to PPTX converter with error handling
    """

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")

    try:
        # Convert PDF pages to images
        print(f"Converting {pdf_path} to images...")
        images = convert_from_path(pdf_path, dpi=dpi)
        print(f"Found {len(images)} pages")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for i, image in enumerate(images):
            print(f"Processing page {i+1}", end=" ... ")

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save image temporarily
            temp_image_path = f"temp_page_{i}_{os.getpid()}.png"
            image.save(temp_image_path, "PNG")

            try:
                # Calculate dimensions to maintain aspect ratio
                img_width, img_height = image.size
                slide_aspect = prs.slide_width / prs.slide_height
                img_aspect = img_width / img_height

                if img_aspect > slide_aspect:
                    # Image is wider - fit to width
                    width = prs.slide_width
                    height = prs.slide_width / img_aspect
                    left = Inches(0)
                    top = (prs.slide_height - height) / 2
                else:
                    # Image is taller - fit to height
                    height = prs.slide_height
                    width = prs.slide_height * img_aspect
                    top = Inches(0)
                    left = (prs.slide_width - width) / 2

                # Add image to slide
                slide.shapes.add_picture(temp_image_path, left, top, width, height)
                print("✓")

            finally:
                # Clean up temporary file
                if os.path.exists(temp_image_path):
                    os.remove(temp_image_path)

        # Save presentation
        prs.save(output_path)
        print(f"\n✅ Successfully converted to {output_path}")

    except Exception as e:
        print(f"\n❌ Error during conversion: {str(e)}")
        sys.exit(1)


# Usage example
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Convert PDF to PPTX")
    parser.add_argument("pdf_file", help="Input PDF file path")
    parser.add_argument("-o", "--output", help="Output PPTX file path")
    parser.add_argument("--dpi", type=int, default=300, help="DPI for image conversion")

    args = parser.parse_args()

    # Set output filename if not provided
    if args.output:
        output_file = args.output
    else:
        base_name = os.path.splitext(args.pdf_file)[0]
        output_file = f"{base_name}.pptx"

    pdf_to_pptx_enhanced(args.pdf_file, output_file, args.dpi)

# usage (windows)
# cd ZJU-GISLAB-COURSE-2025
# conda activate env
# pdf_to_pptx.bat
